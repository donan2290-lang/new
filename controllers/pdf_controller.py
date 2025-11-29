"""
PDF Conversion Controller
Handles all PDF conversion operations:
- PDF to Word, Excel, PowerPoint
- Word/Excel/PowerPoint to PDF
- PDF to Images (JPG, PNG)
- Images to PDF
- Merge PDF, Split PDF, Compress PDF
"""

from flask import Blueprint, request, jsonify, send_file
from werkzeug.utils import secure_filename
import os
import tempfile
from datetime import datetime
import uuid
import time
import shutil
import subprocess

# Import PDF libraries
try:
    from pdf2docx import Converter as PDFToWordConverter
    PDF2DOCX_AVAILABLE = True
except ImportError:
    PDF2DOCX_AVAILABLE = False

try:
    from PIL import Image
    import img2pdf
    IMAGE_LIBS_AVAILABLE = True
except ImportError:
    IMAGE_LIBS_AVAILABLE = False

try:
    from PyPDF2 import PdfMerger, PdfReader, PdfWriter
    PYPDF_AVAILABLE = True
except ImportError:
    try:
        from PyPDF2 import PdfFileMerger as PdfMerger, PdfFileReader as PdfReader, PdfFileWriter as PdfWriter
        PYPDF_AVAILABLE = True
    except ImportError:
        PYPDF_AVAILABLE = False

try:
    from pdf2image import convert_from_path
    PDF2IMAGE_AVAILABLE = True
except ImportError:
    PDF2IMAGE_AVAILABLE = False

try:
    from docx2pdf import convert as docx_to_pdf_convert
    DOCX2PDF_AVAILABLE = True
except ImportError:
    DOCX2PDF_AVAILABLE = False

try:
    import openpyxl
    from openpyxl import load_workbook
    import tabula
    EXCEL_LIBS_AVAILABLE = True
except ImportError:
    EXCEL_LIBS_AVAILABLE = False

try:
    from pptx import Presentation
    from pptx.util import Inches
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

# Create Blueprint
pdf_bp = Blueprint('pdf', __name__, url_prefix='/api/pdf')

# Configuration
UPLOAD_FOLDER = 'uploads'
OUTPUT_FOLDER = 'outputs'
ALLOWED_PDF_EXTENSIONS = {'pdf'}
ALLOWED_IMAGE_EXTENSIONS = {'png', 'jpg', 'jpeg', 'webp'}
ALLOWED_DOC_EXTENSIONS = {'doc', 'docx', 'xlsx','xls', 'pptx'}
MAX_FILE_SIZE = 50 * 1024 * 1024  # 50MB
PDF_TO_WORD_MAX_PAGES = int(os.getenv('PDF_TO_WORD_MAX_PAGES', '0'))  # 0 = no limit

def allowed_file(filename, allowed_extensions):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in allowed_extensions

def get_file_extension(filename):
    return filename.rsplit('.', 1)[1].lower() if '.' in filename else ''

def generate_unique_filename(original_filename, suffix=''):
    """Generate unique filename with timestamp and UUID"""
    name, ext = os.path.splitext(original_filename)
    timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
    unique_id = str(uuid.uuid4())[:8]
    if suffix:
        return f"{name}_{suffix}_{timestamp}_{unique_id}{ext}"
    return f"{name}_{timestamp}_{unique_id}{ext}"

# ==================== PDF to Word ====================
@pdf_bp.route('/to-word', methods=['POST'])
def pdf_to_word():
    """Convert PDF to Word (DOCX)"""
    if not PDF2DOCX_AVAILABLE:
        return jsonify({
            'success': False,
            'error': 'PDF to Word converter not available. Install pdf2docx library.'
        }), 500
    
    try:
        # Check if file exists
        if 'file' not in request.files:
            return jsonify({'success': False, 'error': 'No file provided'}), 400
        
        file = request.files['file']
        if file.filename == '':
            return jsonify({'success': False, 'error': 'No file selected'}), 400
        
        if not allowed_file(file.filename, ALLOWED_PDF_EXTENSIONS):
            return jsonify({'success': False, 'error': 'Invalid file type. Only PDF allowed.'}), 400
        
        # Ensure folders exist to avoid IO errors/delays
        os.makedirs(UPLOAD_FOLDER, exist_ok=True)
        os.makedirs(OUTPUT_FOLDER, exist_ok=True)

        # Save uploaded file
        filename = secure_filename(file.filename)
        pdf_path = os.path.join(UPLOAD_FOLDER, generate_unique_filename(filename))
        file.save(pdf_path)
        
        # Convert PDF to Word
        docx_filename = filename.rsplit('.', 1)[0] + '.docx'
        docx_path = os.path.join(OUTPUT_FOLDER, generate_unique_filename(docx_filename, 'converted'))
        
        # Perform conversion with optional page limit and timing
        start_time = time.time()
        converter = PDFToWordConverter(pdf_path)
        # Determine end page if user supplied param or env var
        end_page = None
        try:
            # Allow client to pass `max_pages` as form field (int)
            max_pages_param = request.form.get('max_pages') or request.args.get('max_pages')
            if max_pages_param:
                mp = int(max_pages_param)
                if mp > 0:
                    end_page = mp
        except Exception:
            end_page = None

        if not end_page and PDF_TO_WORD_MAX_PAGES > 0:
            end_page = PDF_TO_WORD_MAX_PAGES

        # If end_page is set, pass it to converter (pdf2docx uses end as page number limit)
        converter.convert(docx_path, start=0, end=end_page)
        converter.close()
        duration = time.time() - start_time
        print(f"PDF->DOCX conversion finished in {duration:.2f}s (pages limit={end_page})")
        
        # Cleanup uploaded file
        os.remove(pdf_path)
        
        return send_file(
            docx_path,
            as_attachment=True,
            download_name=docx_filename,
            mimetype='application/vnd.openxmlformats-officedocument.wordprocessingml.document'
        )
        
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500

# ==================== Word to PDF ====================
@pdf_bp.route('/from-word', methods=['POST'])
def word_to_pdf():
    """Convert Word (DOCX) to PDF"""
    try:
        if 'file' not in request.files:
            return jsonify({'success': False, 'error': 'No file provided'}), 400

        file = request.files['file']
        if file.filename == '':
            return jsonify({'success': False, 'error': 'No file selected'}), 400

        if not allowed_file(file.filename, ALLOWED_DOC_EXTENSIONS):
            return jsonify({'success': False, 'error': 'Invalid file type. Supported: DOC, DOCX'}), 400

        # Ensure folders exist
        os.makedirs(UPLOAD_FOLDER, exist_ok=True)
        os.makedirs(OUTPUT_FOLDER, exist_ok=True)

        # Save uploaded file
        filename = secure_filename(file.filename)
        docx_path = os.path.join(UPLOAD_FOLDER, generate_unique_filename(filename))
        file.save(docx_path)

        # Output PDF path
        pdf_filename = filename.rsplit('.', 1)[0] + '.pdf'
        pdf_path = os.path.join(OUTPUT_FOLDER, generate_unique_filename(pdf_filename, 'converted'))

        # Try docx2pdf if available (requires MS Word on Windows)
        if DOCX2PDF_AVAILABLE:
            start = time.time()
            docx_to_pdf_convert(docx_path, pdf_path)
            duration = time.time() - start
            print(f"Converted DOCX->PDF using docx2pdf in {duration:.2f}s")

        else:
            # Fallback: try LibreOffice (soffice) CLI if installed
            soffice_path = shutil.which('soffice') or shutil.which('libreoffice')
            if soffice_path:
                # LibreOffice headless conversion
                start = time.time()
                try:
                    # --convert-to pdf --outdir <dir> <file>
                    subprocess.check_call([
                        soffice_path,
                        '--headless',
                        '--convert-to', 'pdf',
                        '--outdir', OUTPUT_FOLDER,
                        docx_path
                    ], stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)

                    # LibreOffice names output file same basename with .pdf
                    converted_name = os.path.splitext(os.path.basename(docx_path))[0] + '.pdf'
                    converted_path = os.path.join(OUTPUT_FOLDER, converted_name)
                    if os.path.exists(converted_path):
                        # Move/rename to our unique pdf_path
                        os.replace(converted_path, pdf_path)
                        duration = time.time() - start
                        print(f"Converted DOCX->PDF using LibreOffice in {duration:.2f}s")
                    else:
                        raise RuntimeError('LibreOffice conversion failed: output not found')
                except subprocess.CalledProcessError as e:
                    raise RuntimeError(f'LibreOffice conversion failed: {e}')
            else:
                # Neither docx2pdf nor libreoffice available
                os.remove(docx_path)
                return jsonify({
                    'success': False,
                    'error': 'Server cannot convert DOCX: install `docx2pdf` (requires MS Word) or LibreOffice (`soffice`).'
                }), 500

        # Cleanup uploaded file
        try:
            if os.path.exists(docx_path):
                os.remove(docx_path)
        except Exception:
            pass

        return send_file(
            pdf_path,
            as_attachment=True,
            download_name=pdf_filename,
            mimetype='application/pdf'
        )

    except Exception as e:
        # Ensure we don't leak uploaded file on error
        try:
            if 'docx_path' in locals() and os.path.exists(docx_path):
                os.remove(docx_path)
        except Exception:
            pass
        return jsonify({'success': False, 'error': str(e)}), 500

# ==================== PDF to JPG ====================
@pdf_bp.route('/to-jpg', methods=['POST'])
def pdf_to_jpg():
    """Convert PDF pages to JPG images using PyMuPDF (fitz)"""
    try:
        if 'file' not in request.files:
            return jsonify({'success': False, 'error': 'No file provided'}), 400
        
        file = request.files['file']
        if not allowed_file(file.filename, ALLOWED_PDF_EXTENSIONS):
            return jsonify({'success': False, 'error': 'Invalid file type'}), 400
        
        # Save PDF
        filename = secure_filename(file.filename)
        pdf_path = os.path.join(UPLOAD_FOLDER, generate_unique_filename(filename))
        file.save(pdf_path)
        
        # Convert to images using PyMuPDF
        import fitz
        doc = fitz.open(pdf_path)
        output_files = []
        
        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            pix = page.get_pixmap(matrix=fitz.Matrix(300/72, 300/72))  # 300 DPI
            
            jpg_filename = f"{filename.rsplit('.', 1)[0]}_page_{page_num+1}.jpg"
            jpg_path = os.path.join(OUTPUT_FOLDER, generate_unique_filename(jpg_filename))
            pix.save(jpg_path)
            output_files.append(jpg_path)
        
        doc.close()
        
        # Cleanup
        os.remove(pdf_path)
        
        # If single page, return file directly
        if len(output_files) == 1:
            return send_file(output_files[0], as_attachment=True, mimetype='image/jpeg')
        
        # Multiple pages - create ZIP
        import zipfile
        zip_path = os.path.join(OUTPUT_FOLDER, f"{filename.rsplit('.', 1)[0]}_images.zip")
        with zipfile.ZipFile(zip_path, 'w') as zipf:
            for file_path in output_files:
                zipf.write(file_path, os.path.basename(file_path))
                os.remove(file_path)
        
        return send_file(zip_path, as_attachment=True, mimetype='application/zip')
        
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500

# ==================== JPG to PDF ====================
@pdf_bp.route('/from-jpg', methods=['POST'])
def jpg_to_pdf():
    """Convert JPG images to PDF"""
    if not IMAGE_LIBS_AVAILABLE:
        return jsonify({
            'success': False,
            'error': 'Image to PDF converter not available. Install img2pdf and Pillow.'
        }), 500
    
    try:
        if 'files' not in request.files:
            return jsonify({'success': False, 'error': 'No files provided'}), 400
        
        files = request.files.getlist('files')
        if not files:
            return jsonify({'success': False, 'error': 'No files selected'}), 400
        
        # Save all images
        image_paths = []
        for file in files:
            if allowed_file(file.filename, ALLOWED_IMAGE_EXTENSIONS):
                filename = secure_filename(file.filename)
                img_path = os.path.join(UPLOAD_FOLDER, generate_unique_filename(filename))
                file.save(img_path)
                image_paths.append(img_path)
        
        if not image_paths:
            return jsonify({'success': False, 'error': 'No valid images provided'}), 400
        
        # Convert to PDF
        pdf_filename = 'images_to_pdf.pdf'
        pdf_path = os.path.join(OUTPUT_FOLDER, generate_unique_filename(pdf_filename, 'converted'))
        
        with open(pdf_path, 'wb') as f:
            f.write(img2pdf.convert(image_paths))
        
        # Cleanup
        for img_path in image_paths:
            os.remove(img_path)
        
        return send_file(pdf_path, as_attachment=True, mimetype='application/pdf')
        
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500

# ==================== PDF to PNG ====================
@pdf_bp.route('/to-png', methods=['POST'])
def pdf_to_png():
    """Convert PDF pages to PNG images using PyMuPDF (fitz)"""
    try:
        if 'file' not in request.files:
            return jsonify({'success': False, 'error': 'No file provided'}), 400
        
        file = request.files['file']
        if not allowed_file(file.filename, ALLOWED_PDF_EXTENSIONS):
            return jsonify({'success': False, 'error': 'Invalid file type'}), 400
        
        filename = secure_filename(file.filename)
        pdf_path = os.path.join(UPLOAD_FOLDER, generate_unique_filename(filename))
        file.save(pdf_path)
        
        # Convert to PNG using PyMuPDF
        import fitz
        doc = fitz.open(pdf_path)
        output_files = []
        
        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            pix = page.get_pixmap(matrix=fitz.Matrix(300/72, 300/72))  # 300 DPI
            
            png_filename = f"{filename.rsplit('.', 1)[0]}_page_{page_num+1}.png"
            png_path = os.path.join(OUTPUT_FOLDER, generate_unique_filename(png_filename))
            pix.save(png_path)
            output_files.append(png_path)
        
        doc.close()
        os.remove(pdf_path)
        
        if len(output_files) == 1:
            return send_file(output_files[0], as_attachment=True, mimetype='image/png')
        
        # Multiple pages - ZIP
        import zipfile
        zip_path = os.path.join(OUTPUT_FOLDER, f"{filename.rsplit('.', 1)[0]}_images.zip")
        with zipfile.ZipFile(zip_path, 'w') as zipf:
            for file_path in output_files:
                zipf.write(file_path, os.path.basename(file_path))
                os.remove(file_path)
        
        return send_file(zip_path, as_attachment=True, mimetype='application/zip')
        
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500

# ==================== PNG to PDF ====================
@pdf_bp.route('/from-png', methods=['POST'])
def png_to_pdf():
    """Convert PNG images to PDF"""
    return jpg_to_pdf()  # Same logic as JPG to PDF

# ==================== Merge PDF ====================
@pdf_bp.route('/merge', methods=['POST'])
def merge_pdf():
    """Merge multiple PDF files into one"""
    if not PYPDF_AVAILABLE:
        return jsonify({
            'success': False,
            'error': 'PDF merger not available. Install PyPDF2.'
        }), 500
    
    try:
        if 'files' not in request.files:
            return jsonify({'success': False, 'error': 'No files provided'}), 400
        
        files = request.files.getlist('files')
        if len(files) < 2:
            return jsonify({'success': False, 'error': 'At least 2 PDF files required'}), 400
        
        # Save all PDFs
        pdf_paths = []
        for file in files:
            if allowed_file(file.filename, ALLOWED_PDF_EXTENSIONS):
                filename = secure_filename(file.filename)
                pdf_path = os.path.join(UPLOAD_FOLDER, generate_unique_filename(filename))
                file.save(pdf_path)
                pdf_paths.append(pdf_path)
        
        if len(pdf_paths) < 2:
            return jsonify({'success': False, 'error': 'At least 2 valid PDF files required'}), 400
        
        # Merge PDFs
        merger = PdfMerger()
        for pdf_path in pdf_paths:
            merger.append(pdf_path)
        
        # Save merged PDF
        output_filename = 'merged_document.pdf'
        output_path = os.path.join(OUTPUT_FOLDER, generate_unique_filename(output_filename, 'merged'))
        merger.write(output_path)
        merger.close()
        
        # Cleanup
        for pdf_path in pdf_paths:
            os.remove(pdf_path)
        
        return send_file(output_path, as_attachment=True, mimetype='application/pdf')
        
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500

# ==================== Split PDF ====================
@pdf_bp.route('/split', methods=['POST'])
def split_pdf():
    """Split PDF into individual pages"""
    if not PYPDF_AVAILABLE:
        return jsonify({
            'success': False,
            'error': 'PDF splitter not available. Install PyPDF2.'
        }), 500
    
    try:
        if 'file' not in request.files:
            return jsonify({'success': False, 'error': 'No file provided'}), 400
        
        file = request.files['file']
        if not allowed_file(file.filename, ALLOWED_PDF_EXTENSIONS):
            return jsonify({'success': False, 'error': 'Invalid file type'}), 400
        
        filename = secure_filename(file.filename)
        pdf_path = os.path.join(UPLOAD_FOLDER, generate_unique_filename(filename))
        file.save(pdf_path)
        
        # Split PDF
        reader = PdfReader(pdf_path)
        output_files = []
        
        for i, page in enumerate(reader.pages):
            writer = PdfWriter()
            writer.add_page(page)
            
            page_filename = f"{filename.rsplit('.', 1)[0]}_page_{i+1}.pdf"
            page_path = os.path.join(OUTPUT_FOLDER, generate_unique_filename(page_filename))
            
            with open(page_path, 'wb') as output_file:
                writer.write(output_file)
            
            output_files.append(page_path)
        
        os.remove(pdf_path)
        
        # Create ZIP with all pages
        import zipfile
        zip_path = os.path.join(OUTPUT_FOLDER, f"{filename.rsplit('.', 1)[0]}_split.zip")
        with zipfile.ZipFile(zip_path, 'w') as zipf:
            for file_path in output_files:
                zipf.write(file_path, os.path.basename(file_path))
                os.remove(file_path)
        
        return send_file(zip_path, as_attachment=True, mimetype='application/zip')
        
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500

# ==================== Compress PDF ====================
@pdf_bp.route('/compress', methods=['POST'])
def compress_pdf():
    """Compress PDF file size"""
    if not PYPDF_AVAILABLE:
        return jsonify({
            'success': False,
            'error': 'PDF compressor not available. Install PyPDF2.'
        }), 500
    
    try:
        if 'file' not in request.files:
            return jsonify({'success': False, 'error': 'No file provided'}), 400
        
        file = request.files['file']
        if not allowed_file(file.filename, ALLOWED_PDF_EXTENSIONS):
            return jsonify({'success': False, 'error': 'Invalid file type'}), 400
        
        filename = secure_filename(file.filename)
        pdf_path = os.path.join(UPLOAD_FOLDER, generate_unique_filename(filename))
        file.save(pdf_path)
        
        # Compress PDF
        reader = PdfReader(pdf_path)
        writer = PdfWriter()
        
        for page in reader.pages:
            page.compress_content_streams()
            writer.add_page(page)
        
        # Save compressed PDF
        compressed_filename = f"{filename.rsplit('.', 1)[0]}_compressed.pdf"
        compressed_path = os.path.join(OUTPUT_FOLDER, generate_unique_filename(compressed_filename))
        
        with open(compressed_path, 'wb') as output_file:
            writer.write(output_file)
        
        os.remove(pdf_path)
        
        # Calculate compression ratio
        original_size = os.path.getsize(pdf_path) if os.path.exists(pdf_path) else file.content_length
        compressed_size = os.path.getsize(compressed_path)
        compression_ratio = ((original_size - compressed_size) / original_size * 100) if original_size > 0 else 0
        
        return send_file(
            compressed_path,
            as_attachment=True,
            download_name=compressed_filename,
            mimetype='application/pdf'
        )
        
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500

# ==================== Excel to PDF ====================
@pdf_bp.route('/from-excel', methods=['POST'])
def excel_to_pdf():
    """Convert Excel to PDF using openpyxl+reportlab or LibreOffice fallback"""
    try:
        if 'file' not in request.files:
            return jsonify({'success': False, 'error': 'No file provided'}), 400

        file = request.files['file']
        if not allowed_file(file.filename, {'xlsx', 'xls'}):
            return jsonify({'success': False, 'error': 'Invalid file type. Only Excel files allowed.'}), 400

        # Ensure folders exist
        os.makedirs(UPLOAD_FOLDER, exist_ok=True)
        os.makedirs(OUTPUT_FOLDER, exist_ok=True)

        filename = secure_filename(file.filename)
        excel_path = os.path.join(UPLOAD_FOLDER, generate_unique_filename(filename))
        file.save(excel_path)

        # Output PDF path
        pdf_filename = filename.rsplit('.', 1)[0] + '.pdf'
        pdf_path = os.path.join(OUTPUT_FOLDER, generate_unique_filename(pdf_filename, 'converted'))

        # Try openpyxl + reportlab if available (supports .xlsx only)
        if EXCEL_LIBS_AVAILABLE and filename.lower().endswith('.xlsx'):
            start = time.time()
            try:
                from reportlab.lib.pagesizes import A4, landscape
                from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph
                from reportlab.lib import colors
                from reportlab.lib.styles import getSampleStyleSheet

                # Load Excel
                wb = load_workbook(excel_path, data_only=True)

                # Create PDF
                doc = SimpleDocTemplate(pdf_path, pagesize=landscape(A4))
                elements = []
                styles = getSampleStyleSheet()

                # Convert each sheet
                for sheet_name in wb.sheetnames:
                    ws = wb[sheet_name]

                    # Add sheet title
                    elements.append(Paragraph(f"<b>{sheet_name}</b>", styles['Heading1']))

                    # Get data
                    data = []
                    for row in ws.iter_rows(values_only=True):
                        data.append([str(cell) if cell is not None else '' for cell in row])

                    if data:
                        # Create table
                        t = Table(data)
                        t.setStyle(TableStyle([
                            ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
                            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                            ('FONTSIZE', (0, 0), (-1, 0), 10),
                            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                            ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                            ('GRID', (0, 0), (-1, -1), 1, colors.black)
                        ]))
                        elements.append(t)

                doc.build(elements)
                duration = time.time() - start
                print(f"Converted XLSX->PDF using openpyxl+reportlab in {duration:.2f}s")

                # Cleanup
                try:
                    if os.path.exists(excel_path):
                        os.remove(excel_path)
                except Exception:
                    pass

                return send_file(pdf_path, as_attachment=True, download_name=pdf_filename, mimetype='application/pdf')

            except Exception as e:
                print(f"openpyxl+reportlab conversion failed: {e}")
                # Fall through to LibreOffice fallback

        # Fallback: try LibreOffice (soffice) CLI if installed
        soffice_path = shutil.which('soffice') or shutil.which('libreoffice')
        if soffice_path:
            start = time.time()
            try:
                # LibreOffice headless conversion
                subprocess.check_call([
                    soffice_path,
                    '--headless',
                    '--convert-to', 'pdf',
                    '--outdir', OUTPUT_FOLDER,
                    excel_path
                ], stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)

                # LibreOffice names output file same basename with .pdf
                converted_name = os.path.splitext(os.path.basename(excel_path))[0] + '.pdf'
                converted_path = os.path.join(OUTPUT_FOLDER, converted_name)
                if os.path.exists(converted_path):
                    # Move/rename to our unique pdf_path
                    os.replace(converted_path, pdf_path)
                    duration = time.time() - start
                    print(f"Converted Excel->PDF using LibreOffice in {duration:.2f}s")
                else:
                    raise RuntimeError('LibreOffice conversion failed: output not found')
            except subprocess.CalledProcessError as e:
                raise RuntimeError(f'LibreOffice conversion failed: {e}')

            # Cleanup
            try:
                if os.path.exists(excel_path):
                    os.remove(excel_path)
            except Exception:
                pass

            return send_file(pdf_path, as_attachment=True, download_name=pdf_filename, mimetype='application/pdf')

        else:
            # Neither method available
            try:
                if os.path.exists(excel_path):
                    os.remove(excel_path)
            except Exception:
                pass
            return jsonify({
                'success': False,
                'error': 'No available converter for Excel files. Install openpyxl (for XLSX) or LibreOffice (soffice) for XLS support.'
            }), 500

    except Exception as e:
        # Ensure we don't leak uploaded file on error
        try:
            if 'excel_path' in locals() and os.path.exists(excel_path):
                os.remove(excel_path)
        except Exception:
            pass
        return jsonify({'success': False, 'error': str(e)}), 500

# ==================== PDF to Excel ====================
@pdf_bp.route('/to-excel', methods=['POST'])
def pdf_to_excel():
    """Convert PDF tables to Excel using PyMuPDF text extraction"""
    try:
        if not EXCEL_LIBS_AVAILABLE:
            return jsonify({
                'success': False,
                'error': 'Excel libraries not available. Install openpyxl.'
            }), 500
        
        if 'file' not in request.files:
            return jsonify({'success': False, 'error': 'No file provided'}), 400
        
        file = request.files['file']
        if not allowed_file(file.filename, ALLOWED_PDF_EXTENSIONS):
            return jsonify({'success': False, 'error': 'Invalid file type. Only PDF allowed.'}), 400
        
        filename = secure_filename(file.filename)
        pdf_path = os.path.join(UPLOAD_FOLDER, generate_unique_filename(filename))
        file.save(pdf_path)
        
        # Extract text using PyMuPDF
        import fitz
        
        doc = fitz.open(pdf_path)
        
        # Create Excel file
        excel_filename = filename.rsplit('.', 1)[0] + '.xlsx'
        excel_path = os.path.join(OUTPUT_FOLDER, generate_unique_filename(excel_filename, 'converted'))
        
        wb = openpyxl.Workbook()
        wb.remove(wb.active)
        
        # Process each page
        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            text = page.get_text()
            
            # Create sheet for page
            sheet_name = f'Page_{page_num+1}'
            ws = wb.create_sheet(sheet_name)
            
            # Split text into lines and write to Excel
            lines = text.split('\n')
            for row_idx, line in enumerate(lines, 1):
                if line.strip():
                    # Try to split by tabs or multiple spaces
                    cells = [cell.strip() for cell in line.split('\t') if cell.strip()]
                    if len(cells) == 1:
                        # Try splitting by multiple spaces
                        import re
                        cells = [cell for cell in re.split(r'\s{2,}', line) if cell.strip()]
                    
                    for col_idx, cell_value in enumerate(cells, 1):
                        ws.cell(row=row_idx, column=col_idx, value=cell_value)
        
        doc.close()
        wb.save(excel_path)
        
        # Cleanup
        os.remove(pdf_path)
        
        return send_file(excel_path, as_attachment=True, 
                        mimetype='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet')
        
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500

# ==================== PowerPoint to PDF ====================
@pdf_bp.route('/from-ppt', methods=['POST'])
def ppt_to_pdf():
    """Convert PowerPoint to PDF using python-pptx and PyMuPDF"""
    try:
        if not PPTX_AVAILABLE:
            return jsonify({
                'success': False,
                'error': 'PowerPoint library not available. Install python-pptx.'
            }), 500
        
        if 'file' not in request.files:
            return jsonify({'success': False, 'error': 'No file provided'}), 400
        
        file = request.files['file']
        if not allowed_file(file.filename, {'pptx'}):
            return jsonify({'success': False, 'error': 'Invalid file type. Only PPTX allowed.'}), 400
        
        filename = secure_filename(file.filename)
        pptx_path = os.path.join(UPLOAD_FOLDER, generate_unique_filename(filename))
        file.save(pptx_path)
        
        # Convert using reportlab
        from reportlab.lib.pagesizes import letter, landscape
        from reportlab.pdfgen import canvas
        from reportlab.lib.utils import ImageReader
        
        # Load presentation
        prs = Presentation(pptx_path)
        
        # Create PDF
        pdf_filename = filename.rsplit('.', 1)[0] + '.pdf'
        pdf_path = os.path.join(OUTPUT_FOLDER, generate_unique_filename(pdf_filename, 'converted'))
        
        c = canvas.Canvas(pdf_path, pagesize=landscape(letter))
        width, height = landscape(letter)
        
        # Convert each slide to image and add to PDF
        temp_images = []
        for i, slide in enumerate(prs.slides):
            # Create simple slide representation
            c.setFont("Helvetica-Bold", 24)
            c.drawString(50, height - 100, f"Slide {i+1}")
            
            y_position = height - 150
            c.setFont("Helvetica", 14)
            
            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text = shape.text[:100]  # Limit text length
                    c.drawString(50, y_position, text)
                    y_position -= 30
                    if y_position < 100:
                        break
            
            c.showPage()
        
        c.save()
        
        # Cleanup
        os.remove(pptx_path)
        for img in temp_images:
            if os.path.exists(img):
                os.remove(img)
        
        return send_file(pdf_path, as_attachment=True, mimetype='application/pdf')
        
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500

# ==================== PDF to PowerPoint ====================
@pdf_bp.route('/to-ppt', methods=['POST'])
def pdf_to_ppt():
    """Convert PDF to PowerPoint by converting each page to image and inserting into slides"""
    try:
        if not PPTX_AVAILABLE:
            return jsonify({
                'success': False,
                'error': 'PowerPoint library not available. Install python-pptx.'
            }), 500
        
        if 'file' not in request.files:
            return jsonify({'success': False, 'error': 'No file provided'}), 400
        
        file = request.files['file']
        if not allowed_file(file.filename, ALLOWED_PDF_EXTENSIONS):
            return jsonify({'success': False, 'error': 'Invalid file type. Only PDF allowed.'}), 400
        
        filename = secure_filename(file.filename)
        pdf_path = os.path.join(UPLOAD_FOLDER, generate_unique_filename(filename))
        file.save(pdf_path)
        
        # Convert PDF pages to images
        import fitz
        from pptx.util import Inches
        
        doc = fitz.open(pdf_path)
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        temp_images = []
        
        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))  # 2x scale for quality
            
            # Save temporary image
            temp_img = os.path.join(OUTPUT_FOLDER, f"temp_slide_{page_num}.png")
            pix.save(temp_img)
            temp_images.append(temp_img)
            
            # Add slide with image
            blank_slide_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # Add image to slide
            left = Inches(0)
            top = Inches(0)
            slide.shapes.add_picture(temp_img, left, top, width=prs.slide_width, height=prs.slide_height)
        
        doc.close()
        
        # Save PowerPoint
        pptx_filename = filename.rsplit('.', 1)[0] + '.pptx'
        pptx_path = os.path.join(OUTPUT_FOLDER, generate_unique_filename(pptx_filename, 'converted'))
        prs.save(pptx_path)
        
        # Cleanup
        os.remove(pdf_path)
        for img in temp_images:
            if os.path.exists(img):
                os.remove(img)
        
        return send_file(pptx_path, as_attachment=True,
                        mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation')
        
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500

# ==================== Health Check ====================
@pdf_bp.route('/health', methods=['GET'])
def health_check():
    """Check which PDF libraries are available"""
    return jsonify({
        'success': True,
        'libraries': {
            'pdf2docx': PDF2DOCX_AVAILABLE,
            'docx2pdf': DOCX2PDF_AVAILABLE,
            'pypdf': PYPDF_AVAILABLE,
            'pdf2image': PDF2IMAGE_AVAILABLE,
            'pillow': IMAGE_LIBS_AVAILABLE,
            'img2pdf': IMAGE_LIBS_AVAILABLE,
            'excel': EXCEL_LIBS_AVAILABLE,
            'pptx': PPTX_AVAILABLE
        }
    })

# ==================== Generic Image Endpoints (aliases for backward compatibility) ====================
@pdf_bp.route('/to-image', methods=['POST'])
def pdf_to_image():
    """Alias endpoint that routes to either /to-jpg or /to-png based on 'format' parameter.
    Falls back to JPG if not specified.
    """
    fmt = (request.args.get('format') or request.form.get('format') or request.headers.get('X-Image-Format') or '').lower()
    if fmt.startswith('png'):
        return pdf_to_png()
    # default to jpg
    return pdf_to_jpg()


@pdf_bp.route('/from-image', methods=['POST'])
def image_to_pdf():
    """Alias endpoint that routes to either /from-jpg or /from-png based on 'format' parameter.
    Falls back to JPG if not specified.
    """
    fmt = (request.args.get('format') or request.form.get('format') or request.headers.get('X-Image-Format') or '').lower()
    if fmt.startswith('png'):
        return png_to_pdf()
    return jpg_to_pdf()
